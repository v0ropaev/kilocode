#!/usr/bin/env python3
# Builds submission/Security-Auto-ATH.pptx from the AI Talent Hub template.
# Starts from a copy of the template so theme, master, layouts and the embedded Inter
# font travel with the file; the knowbase slides are removed and replaced.

import os
import shutil
import sys
import tempfile
import zipfile
from datetime import datetime, timezone

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

HERE = os.path.dirname(os.path.abspath(__file__))

# The AI Talent Hub template is the visual source: slide size, theme, master, layouts, the embedded
# Inter font and the logo images all travel with it. It is not vendored into the repository (it is a
# different project's deck and carries third-party marks), so its path is given explicitly.
TEMPLATE = os.environ.get("ATH_TEMPLATE") or os.path.join(
    os.path.dirname(os.path.dirname(HERE)), "knowbase-\u043f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u044f-ATH.pptx"
)
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "Security-Auto-ATH.pptx")

# Logo images, by their entry name inside the template package.
MEDIA_ENTRIES = {
    "ath_dark": "ppt/media/image3.png",
    "ath_light": "ppt/media/image5.png",
    "itmo_dark": "ppt/media/image2.png",
    "itmo_light": "ppt/media/image9.png",
    "partner_dark": "ppt/media/image1.png",
    "partner_light": "ppt/media/image10.png",
}
_MEDIA_DIR = None


def media(key):
    """Extract a logo from the template package on first use and return its path."""
    global _MEDIA_DIR
    if _MEDIA_DIR is None:
        _MEDIA_DIR = tempfile.mkdtemp(prefix="ath-media-")
        with zipfile.ZipFile(TEMPLATE) as z:
            names = set(z.namelist())
            missing = [e for e in MEDIA_ENTRIES.values() if e not in names]
            if missing:
                raise SystemExit(
                    f"{TEMPLATE} does not contain the expected logo entries: {', '.join(missing)}.\n"
                    "This build script is tied to the AI Talent Hub deck it was written against."
                )
            for entry in MEDIA_ENTRIES.values():
                target = os.path.join(_MEDIA_DIR, os.path.basename(entry))
                with open(target, "wb") as fh:
                    fh.write(z.read(entry))
    return os.path.join(_MEDIA_DIR, os.path.basename(MEDIA_ENTRIES[key]))


def preflight():
    if not os.path.isfile(TEMPLATE):
        raise SystemExit(
            f"AI Talent Hub template not found at {TEMPLATE}.\n"
            "Pass it explicitly:  ATH_TEMPLATE=/path/to/ath-template.pptx python3 build-deck.py out.pptx"
        )


# ---------------------------------------------------------------- design tokens
BLACK = RGBColor(0x00, 0x00, 0x00)
BLUE = RGBColor(0x00, 0x00, 0xFF)
PURPLE = RGBColor(0x88, 0x43, 0xFF)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = "FAFAFC"
BG_BLUE = "0000FF"

INTER, ARIAL = "Inter", "Arial"

M = 0.521                     # left margin
RIGHT = 9.479                 # right edge
LOGO = (0.521, 0.351, 1.68, 0.464)
LOGO_ON_BLUE = (0.330, 0.336, 1.83, 0.506)
YEAR_POS = (8.799, 0.456, 0.68, 0.293)
TITLE_Y = 0.974
LEDE_Y = 1.640

COL4_X = [0.521, 2.881, 5.242, 7.602]
COL4_W = 1.875
COL3_X = [0.521, 3.633, 6.748]
COL3_W = 2.726
COL2_X = [0.521, 5.250]
COL2_W = 4.229


YEAR = "2026"

# ---------------------------------------------------------------- layout checker
# LibreOffice cannot read the template's embedded Inter and falls back to Arial, which is wider than
# Inter at every weight. Measuring against that substitution is the pessimistic case: what fits in the
# rendered PDF also fits in PowerPoint, where the embedded font is used.
from PIL import ImageFont  # noqa: E402

_FONTDIR = "/System/Library/Fonts/Supplemental/"
_SUBST = {
    ("Inter", True): _FONTDIR + "Arial Bold.ttf",
    ("Inter", False): _FONTDIR + "Arial.ttf",
    ("Arial", True): _FONTDIR + "Arial Bold.ttf",
    ("Arial", False): _FONTDIR + "Arial.ttf",
}
_cache = {}
BOXES = []          # (slide_no, kind, x, y, w, declared_h, needed_h, preview)
CURRENT = [0]
LB = [0.0]          # bottom of the last placed block


def below(gap=0.16):
    return LB[0] + gap


def _font(name, bold, size_pt):
    key = (name, bold, round(size_pt, 1))
    if key not in _cache:
        _cache[key] = ImageFont.truetype(_SUBST[(name, bold)], max(int(round(size_pt * 96 / 72)), 6))
    return _cache[key]


# A rendered line is taller than size*line: the renderer multiplies the font's own line height
# (ascent + descent + gap, ~1.17 em for Arial) by the spacing factor. And it breaks lines slightly
# earlier than PIL's advance widths suggest — measurably so for the bold display face. Both are
# compensated here so the check stays pessimistic.
LINE_FACTOR = 1.17
WIDTH_FACTOR = 0.92


def measure(txt, size_pt, name, bold, width_in, line=1.2, space_pt=0):
    """Approximate rendered height in inches, wrapping at width_in."""
    f = _font(name, bold, size_pt)
    limit = width_in * 96 * WIDTH_FACTOR
    total = 0.0
    for hard in txt.split("\n"):
        words, cur, n = hard.split(), "", 0
        for w in words:
            trial = (cur + " " + w).strip()
            if f.getlength(trial) <= limit or not cur:
                cur = trial
            else:
                n += 1
                cur = w
        n += 1
        total += n * size_pt * line * LINE_FACTOR / 72.0 + space_pt / 72.0
    return total


def note(x, y, w, declared_h, needed_h, preview, kind="text"):
    BOXES.append((CURRENT[0], kind, x, y, w, declared_h, needed_h, preview[:44]))


def report():
    print("\n--- layout check (pessimistic Arial/Arial Black substitution) ---")
    bad = 0
    for sl, kind, x, y, w, dh, nh, prev in BOXES:
        if nh > dh + 0.005:
            bad += 1
            print(f"  OVERFLOW slide {sl:2} {kind:6} y={y:.3f} w={w:.3f} declared={dh:.3f} needed={nh:.3f}  {prev!r}")
    for sl, kind, x, y, w, dh, nh, prev in BOXES:
        if y + max(dh, nh) > 5.50:
            bad += 1
            print(f"  OFFSLIDE slide {sl:2} {kind:6} y={y:.3f} bottom={y + max(dh, nh):.3f}  {prev!r}")
    # vertical collisions between horizontally overlapping boxes on the same slide
    by_slide = {}
    for b in BOXES:
        by_slide.setdefault(b[0], []).append(b)
    for sl, items in sorted(by_slide.items()):
        items = sorted(items, key=lambda b: b[3])
        for i, a in enumerate(items):
            for b in items[i + 1:]:
                ax0, ax1 = a[2], a[2] + a[4]
                bx0, bx1 = b[2], b[2] + b[4]
                if min(ax1, bx1) - max(ax0, bx0) <= 0.02:
                    continue
                ay1 = a[3] + a[6]
                if b[3] < ay1 - 0.012:
                    bad += 1
                    print(f"  COLLIDE  slide {sl:2} y={a[3]:.3f}+{a[6]:.3f} over y={b[3]:.3f}"
                          f"  {a[7]!r} / {b[7]!r}")
    print(f"--- {bad} layout problem(s) ---")
    return bad


# ---------------------------------------------------------------- primitives


def set_bg(slide, hexcolor):
    """Paint the slide background.

    `<p:bg>` has to be the FIRST CHILD OF `<p:cSld>`. Hung anywhere else it is schema-invalid, and
    the renderers disagree about what to do with that: LibreOffice paints it anyway, PowerPoint and
    Keynote drop it silently. The blue slides carry white text, so a dropped fill turns them into
    blank white pages — and checking the PDF, which LibreOffice produces, cannot catch it. `main()`
    asserts the placement for exactly that reason."""
    from lxml import etree

    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    pns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    csld = slide._element.find(f"{{{pns}}}cSld")
    for old in list(slide._element.findall(f"{{{pns}}}bg")) + list(csld.findall(f"{{{pns}}}bg")):
        old.getparent().remove(old)
    bg = etree.Element(f"{{{pns}}}bg")
    bgpr = etree.SubElement(bg, f"{{{pns}}}bgPr")
    fill = etree.SubElement(bgpr, f"{{{ns}}}solidFill")
    clr = etree.SubElement(fill, f"{{{ns}}}srgbClr")
    clr.set("val", hexcolor)
    etree.SubElement(bgpr, f"{{{ns}}}effectLst")
    csld.insert(0, bg)


def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return box


def para(tf, text, size, bold=False, color=BLACK, font=ARIAL, space=0, line=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return p


def text(slide, x, y, w, h, lines, size=12, bold=False, color=BLACK, font=ARIAL, space=2, line=None,
         align=None):
    """lines: str or list of str (or (text, overrides) tuples)."""
    if isinstance(lines, str):
        lines = [lines]
    box = tb(slide, x, y, w, h)
    tf = box.text_frame
    needed = 0.0
    preview = ""
    for i, ln in enumerate(lines):
        over = {}
        if isinstance(ln, tuple):
            ln, over = ln
        p = para(
            tf,
            ln,
            over.get("size", size),
            over.get("bold", bold),
            over.get("color", color),
            over.get("font", font),
            over.get("space", space),
            over.get("line", line),
            first=(i == 0),
        )
        if align:
            p.alignment = align
        needed += measure(ln, over.get("size", size), over.get("font", font), over.get("bold", bold),
                          w, over.get("line", line) or 1.2, over.get("space", space))
        preview = preview or ln
    note(x, y, w, h, needed, preview)
    LB[0] = y + max(h, needed)
    return box


def rect(slide, x, y, w, h, fill=None, lineclr=None, linew=0.75, shape=MSO_SHAPE.RECTANGLE,
         quiet=False):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    # Drop the theme style reference the autoshape is born with: it carries an outline glow that the
    # renderer paints around the shape's own text. Fill, line and text colour are set explicitly below.
    style = s._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        s._element.remove(style)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if lineclr is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = lineclr
        s.line.width = Pt(linew)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    s.text_frame.margin_left = s.text_frame.margin_right = Inches(0.06)
    s.text_frame.margin_top = s.text_frame.margin_bottom = Inches(0.04)
    s.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    if not quiet:
        note(x, y, w, h, h, "<shape>", kind="shape")
    return s


def hline(slide, x, y, w, color=BLACK, width=0.75):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def arrow(slide, x, y, w, color=GRAY):
    """A slim right-pointing arrow between flow boxes."""
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y - 0.045), Inches(w), Inches(0.09))
    style = a._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
    if style is not None:
        a._element.remove(style)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def pic(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


# ---------------------------------------------------------------- slide chrome


def new_slide(prs, blue=False, year=True, logo=True):
    slide = prs.slides.add_slide(prs.slide_layouts[11])  # DEFAULT: no placeholders
    CURRENT[0] += 1
    set_bg(slide, BG_BLUE if blue else BG)
    if logo:
        if blue:
            pic(slide, media("ath_light"), *LOGO_ON_BLUE)
        else:
            pic(slide, media("ath_dark"), *LOGO)
    if year:
        b = tb(slide, *YEAR_POS)
        p = para(b.text_frame, YEAR, 17, False, WHITE if blue else BLACK, INTER, first=True)
        p.alignment = PP_ALIGN.RIGHT
    return slide


def title(slide, txt, blue=False, y=TITLE_Y, w=8.2, size=26):
    h = max(measure(txt, size, INTER, True, w, 1.12), 0.42)
    text(slide, M, y, w, h, txt, size=size, bold=True,
         color=WHITE if blue else BLACK, font=INTER, line=1.12, space=0)
    return y + h


def lede(slide, txt, y=LEDE_Y, w=8.4, blue=False):
    h = max(measure(txt, 13, INTER, True, w, 1.2), 0.24)
    text(slide, M, y, w, h, txt, size=13, bold=True,
         color=WHITE if blue else BLACK, font=INTER, space=0)
    return y + h


def stat_row(slide, items, y, xs=None, w=None, numsize=26, blue=False, capsize=12):
    """items: [(number, caption)]"""
    xs = xs or (COL4_X if len(items) == 4 else [M + i * (8.958 / len(items)) for i in range(len(items))])
    w = w or (COL4_W if len(items) == 4 else 8.958 / len(items) - 0.4)
    bottom = y
    for (num, cap), x in zip(items, xs):
        nh = numsize * 1.2 * LINE_FACTOR / 72 + 0.03
        text(slide, x, y, w, nh, num, size=numsize, bold=True,
             color=WHITE if blue else BLUE, font=INTER, space=0)
        ch = max(measure(cap, capsize, ARIAL, False, w, 1.15), 0.22)
        text(slide, x, y + nh + 0.10, w, ch, cap, size=capsize,
             color=WHITE if blue else BLACK, font=ARIAL, line=1.15, space=0)
        bottom = max(bottom, y + nh + 0.10 + ch)
    LB[0] = bottom
    return bottom


def cards(slide, items, y, xs, w, numcolor=BLUE, body_h=0.95, gap=0.232, size=12, foot_gap=0.10,
          footcolor=BLUE, footbold=True, footsize=None):
    """items: [(label, body)] or [(label, body, foot)]. Labels, bodies and feet are aligned across
    the row, so a longer column never drags its neighbours' baselines out of line."""
    lh = max([measure(i[0], size, INTER, False, w, 1.2) for i in items] + [0.20])
    yy = y + max(gap, lh + 0.03)
    bh = max([measure(i[1], size, ARIAL, False, w, 1.2) for i in items] + [body_h])
    fs = footsize or (size - 1)
    fh = max([measure(i[2], fs, INTER, footbold, w, 1.25) for i in items if len(i) > 2] + [0.0])
    for item, x in zip(items, xs):
        text(slide, x, y, w, lh, item[0], size=size, color=numcolor, font=INTER, space=0)
        text(slide, x, yy, w, bh, item[1], size=size, font=ARIAL, line=1.2, space=0)
        if len(item) > 2:
            text(slide, x, yy + bh + foot_gap, w, fh, item[2], size=fs, bold=footbold,
                 color=footcolor, font=INTER, line=1.25, space=0)
    LB[0] = yy + bh + (foot_gap + fh if fh else 0)
    return LB[0]


def bullets(slide, items, y, x=M, w=8.4, pitch=0.62, color=PURPLE):
    for i, it in enumerate(items):
        yy = y + i * pitch
        text(slide, x - 0.166, yy - 0.004, 0.10, 0.224, "•", size=12, bold=True, color=color, font=ARIAL)
        text(slide, x, yy, w, pitch - 0.06, it, size=12, font=ARIAL, line=1.2)


def footer_logos(slide, blue=False):
    if blue:
        pic(slide, media("itmo_light"), 5.884, 4.895, 0.774, 0.209)
        pic(slide, media("partner_light"), 7.647, 4.959, 1.832, 0.150)
    else:
        pic(slide, media("itmo_dark"), 0.526, 4.896, 0.774, 0.209)
        pic(slide, media("partner_dark"), 2.288, 4.954, 1.832, 0.151)




# ---------------------------------------------------------------- extra primitives

TRACK = RGBColor(0xDD, 0xDD, 0xE3)
UI_DIR = os.path.join(HERE, "evidence", "ui")


def shot(slide, name, x, y, w, caption=None):
    """Place a recorded terminal session. The height follows the image's own aspect ratio, so a
    longer transcript can never silently overlap whatever comes next."""
    path = os.path.join(UI_DIR, name)
    if not os.path.isfile(path):
        raise SystemExit(
            f"interface capture {path} is missing.\n"
            "Record and render them first:  bash submission/evidence/ui/capture-all.sh "
            "&& bash submission/evidence/ui/render-all.sh"
        )
    top = y
    if caption:
        ch = max(measure(caption, 11, INTER, False, w, 1.2), 0.19)
        text(slide, x, y, w, ch, caption, size=11, color=BLUE, font=INTER, space=0)
        top = y + ch + 0.06
    shape = slide.shapes.add_picture(path, Inches(x), Inches(top), width=Inches(w))
    h = shape.height / 914400.0
    note(x, top, w, h, h, f"<{name}>", kind="image")
    LB[0] = top + h
    return LB[0]


def scorecard(slide, rows, y, namew=3.40, trackw=1.95, pitch=0.281):
    """One line per kind of threat: how many attack runs reached real damage without the layer
    (the whole grey track) and how many still do with it (the blue fill)."""
    bx = M + namew + 0.12
    for i, (name, done, total, status) in enumerate(rows):
        yy = y + i * pitch
        text(slide, M, yy - 0.020, namew, 0.20, name, size=10, font=ARIAL, space=0)
        rect(slide, bx, yy, trackw, 0.125, fill=TRACK, quiet=True)
        w = trackw * done / float(total)
        if w > 0.005:
            rect(slide, bx, yy, max(w, 0.028), 0.125, fill=BLUE, quiet=True)
        text(slide, bx + trackw + 0.14, yy - 0.030, 0.98, 0.20, f"{done}/{total}", size=10, bold=True,
             font=ARIAL, color=BLUE if done else BLACK, space=0)
        text(slide, bx + trackw + 1.20, yy - 0.030, 1.72, 0.20, status, size=10, font=ARIAL,
             color=GRAY, space=0)
    LB[0] = y + len(rows) * pitch
    return LB[0]


def vchain(slide, steps, x, y, w, bh=0.56, gap=0.11):
    """A top-down chain: each step narrows what the one above it allowed."""
    for i, (head, sub) in enumerate(steps):
        yy = y + i * (bh + gap)
        r = rect(slide, x, yy, w, bh, fill=BLUE if i == 0 else None, lineclr=None if i == 0 else BLACK)
        tf = r.text_frame
        para(tf, head, 10.5, True, WHITE if i == 0 else BLACK, INTER, space=1, first=True)
        para(tf, sub, 8.5, False, WHITE if i == 0 else GRAY, ARIAL, space=0)
        if i < len(steps) - 1:
            a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x + w / 2 - 0.05),
                                       Inches(yy + bh + 0.025), Inches(0.10), Inches(gap - 0.05))
            style = a._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
            if style is not None:
                a._element.remove(style)
            a.fill.solid()
            a.fill.fore_color.rgb = GRAY
            a.line.fill.background()
            a.shadow.inherit = False
    LB[0] = y + len(steps) * (bh + gap) - gap
    return LB[0]


def ladder(slide, y, x=M, w=8.958, h=1.05):
    """The ten measured configurations, one bar each: every step switches on exactly one layer."""
    vals = [100, 71, 60, 57, 50, 43, 38, 31, 24, 5]
    bw = (w - 9 * 0.07) / 10.0
    for i, v in enumerate(vals):
        bh = max(h * v / 100.0, 0.035)
        xx = x + i * (bw + 0.07)
        rect(slide, xx, y + h - bh, bw, bh, fill=BLUE if i == len(vals) - 1 else TRACK, quiet=True)
        lab = text(slide, xx, y + h - bh - 0.20, bw, 0.18, f"{v} %", size=9,
                   bold=(i == len(vals) - 1), color=BLUE if i == len(vals) - 1 else GRAY,
                   font=ARIAL, space=0)
        lab.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text(slide, x, y + h + 0.06, 1.8, 0.20, "без защиты · 100 %", size=9, color=GRAY, font=ARIAL, space=0)
    p = text(slide, x + w - 2.4, y + h + 0.06, 2.4, 0.20, "все слои включены · 5 %", size=9, bold=True,
             color=BLUE, font=ARIAL, space=0)
    p.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    LB[0] = y + h + 0.26
    return LB[0]


def statement(slide, y, head, body, w=8.958):
    """A boxed conclusion: the sentence the slide exists to make."""
    hh = 0.26 + measure(head, 13, INTER, True, w - 0.46, 1.2) + 0.06 + \
        measure(body, 11, ARIAL, False, w - 0.46, 1.25) + 0.20
    r = rect(slide, M, y, w, hh, fill=None, lineclr=BLACK, linew=0.75)
    tf = r.text_frame
    tf.margin_left = tf.margin_right = Inches(0.23)
    tf.margin_top = Inches(0.13)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    para(tf, head, 13, True, BLACK, INTER, space=4, first=True)
    para(tf, body, 11, False, BLACK, ARIAL, space=0, line=1.25)
    LB[0] = y + hh
    return LB[0]


def qa_list(slide, items, x, y, w, pitch=0.86):
    """Question, one-line answer, and the internal name of the layer in small grey type."""
    for i, (q, a, tag) in enumerate(items):
        yy = y + i * pitch
        text(slide, x, yy, w, 0.21, q, size=11, color=BLUE, font=INTER, space=0)
        ah = max(measure(a, 10, ARIAL, False, w, 1.2), 0.20)
        text(slide, x, yy + 0.24, w, ah, a, size=10, font=ARIAL, line=1.2, space=0)
        text(slide, x, yy + 0.27 + ah, w, 0.17, tag, size=8.5, color=GRAY, font=ARIAL, space=0)
    LB[0] = y + len(items) * pitch
    return LB[0]


# ---------------------------------------------------------------- content


def pairs(slide, items, y, labelw, size=20, subsize=None, gap=0.15, labelcolor=BLUE,
          subcolor=BLACK, blue=False):
    """Rows of `label | body`, both on one line where the text allows.

    The single layout that survives the organiser's font floor. A three-column grid at 22 pt holds
    about fifteen characters per line — enough for a heading and nothing else — so prose has to run
    across the slide, and the only vertical structure left is the row."""
    bodyx = M + labelw + 0.11
    bodyw = 8.958 - labelw - 0.11
    ss = subsize or size
    for label, body in items:
        lh = max(measure(label, size, INTER, True, labelw, 1.2), 0.20)
        bh = max(measure(body, ss, ARIAL, False, bodyw, 1.2), 0.20)
        h = max(lh, bh)
        text(slide, M, y, labelw, h, label, size=size, bold=True,
             color=WHITE if blue else labelcolor, font=INTER, space=0)
        text(slide, bodyx, y, bodyw, h, body, size=ss,
             color=WHITE if blue else subcolor, font=ARIAL, line=1.2, space=0)
        y += h + gap
    LB[0] = y - gap
    return LB[0]


def build(prs):
    # =============================================================== MAIN DECK — seven slides
    # The organiser's brief: five to seven content slides, one thought each, a title that states a
    # conclusion, and body type a back row can read. Everything that explains rather than convinces
    # is in the appendix, behind the divider.

    # ---------------------------------------------------------- 1. what we built  (0:00–0:30)
    s = new_slide(prs, year=False)
    y = title(s, "Агент работает сам — прав\nпользователя не получает", size=30, w=8.958)
    y = text(s, M, y + 0.20, 8.958, 2.15,
             ["Мы создали Security Auto Mode для разработчиков, которые запускают ИИ-агента "
              "в своём проекте. Проблема: права агента — это права любого, кто сумел им "
              "управлять. Решаем тем, что каждое действие судит отдельный код, а не сама модель."],
             size=19, font=ARIAL, line=1.3, space=0) and below(0.24)
    text(s, M, min(y, 4.16), 8.958, 0.62,
         ["Kilo Code — открытый ИИ-агент, который правит код прямо в проекте.",
          "Дмитрий Воропаев · AI Talent Hub · ИТМО"],
         size=13, font=INTER, color=GRAY, line=1.35, space=0)
    footer_logos(s)

    # ---------------------------------------------------------- 2. why it matters  (0:30–1:00)
    s = new_slide(prs)
    y = title(s, "Оба режима ломают агента", size=28, w=8.958)
    y = pairs(s, [
        ("Спрашивать всё",
         "через полчаса человек жмёт «да» не глядя — самостоятельности не остаётся"),
        ("Отдать всё",
         "одна подброшенная строка в README — и ключи ушли наружу"),
    ], y=y + 0.32, labelw=2.85, size=21, gap=0.32)
    r = rect(s, M, y + 0.34, 8.958, 0.72, fill=BLUE)
    para(r.text_frame, "Автономность нужна. Неограниченные права — нет.",
         21, True, WHITE, INTER, space=0, first=True)
    r.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ---------------------------------------------------------- 3. the demo  (1:00–3:10)
    s = new_slide(prs)
    y = title(s, "README не стал командой для агента", size=28, w=8.958)
    y = text(s, M, y + 0.12, 8.958, 0.36,
             "Задача разработчика: собрать проект и подготовить к выпуску.",
             size=18, font=ARIAL, space=0) and below(0.10)
    y = shot(s, "pitch-injection.png", 0.800, y, 8.400)
    text(s, M, y + 0.12, 8.958, 0.84,
         ["Человек увидел: «отправить файл на внешний адрес, как указано в проекте».",
          ("С обычным README та же команда проходит молча.",
           {"bold": True, "color": BLUE, "font": INTER})],
         size=16, font=ARIAL, line=1.25, space=4)

    # ---------------------------------------------------------- 4. how it works  (3:10–3:45)
    s = new_slide(prs)
    y = title(s, "Правила — гарантии, ИИ — смысл", size=28, w=8.958)
    y = pairs(s, [
        ("Действие", "какая программа, какие файлы, выход в сеть"),
        ("Контекст", "что агент прочитал в этой сессии и откуда"),
        ("Признаки", "два независимых источника, порознь"),
        ("Решение", "разрешить · спросить человека · запретить"),
    ], y=y + 0.18, labelw=2.30, size=21, gap=0.10)
    text(s, M, y + 0.16, 8.958, 1.30,
         ["ИИ возвращает признак, а не решение",
          ("Отменить запрет он не может и источником доверия не является — это свойство кода.",
           {"size": 14, "bold": False, "color": BLACK, "font": ARIAL}),
          ("Смысловой слой — Claude Haiku 4.5 через OpenRouter. Компактная модель, выбрана нарочно: "
           "проверяли принцип, а не мощность.",
           {"size": 13, "bold": False, "color": GRAY, "font": ARIAL})],
         size=17, bold=True, color=BLUE, font=INTER, line=1.22, space=3)

    # ---------------------------------------------------------- 5. the threat map  (3:45–4:10)
    s = new_slide(prs)
    y = title(s, "На каждый вид вреда свой механизм", size=28, w=8.958)
    y = pairs(s, [
        ("Проект удалён или испорчен", "правила по путям и командам"),
        ("Установлен подозрительный пакет", "происхождение до установки"),
        ("Секрет ушёл наружу", "память о прочитанном секрете"),
        ("В README подброшена инструкция", "смысловой слой на модели"),
        ("Инструмент делает лишнее", "права по происхождению"),
        ("Код проекта запускается сам", "запуск после одобрения"),
        ("Расширение читает чужие файлы", "урезанный процесс ОС"),
    ], y=y + 0.16, labelw=5.30, size=19, subsize=15, gap=0.03,
        labelcolor=BLACK, subcolor=GRAY)
    text(s, M, y + 0.12, 8.958, 0.92,
         ["До установки проверяем существование, возраст, распространённость, репозиторий и "
          "install-скрипты — код пакета ещё не запускался.",
          ("slopsquatting: агент выдумывает имя, атакующий его регистрирует · 42 из 42 → 0 из 42",
           {"bold": True, "color": BLUE, "font": INTER})],
         size=13, font=ARIAL, line=1.22, space=3)

    # ---------------------------------------------------------- 6. what we proved  (4:10–4:40)
    s = new_slide(prs, blue=True)
    y = title(s, "Из ста атак доходит пять", blue=True, size=28, w=8.958)
    y = text(s, M, y + 0.20, 8.958, 0.78, "100 %  →  24 %  →  5 %",
             size=40, bold=True, color=WHITE, font=INTER, space=0) and below(0.06)
    y = text(s, M, y, 8.958, 0.34,
             "без защиты  →  одни правила  →  правила и смысловой слой",
             size=16, color=WHITE, font=ARIAL, space=0) and below(0.24)
    y = stat_row(s, [
        ("98 %", "обычных задач выполнено"),
        ("0", "безопасных действий запрещено"),
        ("5760", "прогонов, ни одного сорванного"),
    ], y=y, numsize=26, blue=True, capsize=15)
    text(s, M, y + 0.14, 8.958, 0.88,
         "Успешность атак — доля прогонов, дошедших до наблюдаемого следа: файл добрался до "
         "приёмника, проект исчез. Не по словам модели. 417 прогонов атак, изолированный стенд.",
         size=14, color=WHITE, font=ARIAL, line=1.25, space=0)

    # ---------------------------------------------------------- 7. done / verified / next  (4:40–5:00)
    s = new_slide(prs)
    y = title(s, "Работает, измерено, дальше — пилот", size=28, w=8.958)
    y = pairs(s, [
        ("Сделано", "встроено в Kilo · пакеты · смысловой слой · процесс ОС"),
        ("Проверено", "5760 прогонов · настоящая модель · 755 тестов"),
        ("Дальше", "пилот · два недостающих правила · трение на людях"),
    ], y=y + 0.22, labelw=2.10, size=20, gap=0.14)
    r = rect(s, M, y + 0.16, 8.958, 1.30, fill=BLUE)
    tf = r.text_frame
    tf.margin_left = tf.margin_right = Inches(0.26)
    para(tf, "Security Auto Mode снижает успешность атак со ста процентов до пяти, сохраняя "
             "девяносто восемь процентов обычных задач: агент остаётся автономным, но больше "
             "не получает неограниченные права пользователя.",
         15, True, WHITE, INTER, space=0, line=1.3, first=True)

    # ---------------------------------------------------------- 8. team
    s = new_slide(prs, blue=True)
    y = title(s, "Команда и распределение задач", blue=True, y=0.980)
    hdr = ["Участник", "Роль", "За что отвечал", "Чем подтверждается", "Участие"]
    xs = [0.521, 2.221, 3.921, 6.221, 8.421]
    ws = [1.620, 1.620, 2.220, 2.120, 1.058]
    hy = y + 0.28
    for x, w, h in zip(xs, ws, hdr):
        text(s, x, hy, w, 0.22, h, size=10, bold=True, color=WHITE, font=INTER, space=0)
    hline(s, M, hy + 0.28, 8.958, WHITE, 0.75)
    rows = [
        ["Дмитрий Воропаев", "Разработка,\nAI Product",
         "Архитектура движка решений, слои, процесс расширений, измерение, документация",
         "77 коммитов из 95 в ветке, видно в истории репозитория", "ведущая"],
        ["Артур Фахретдинов", "Разработка",
         "Расширение корпуса атак, два канала egress в разборе curl, первый прототип слоя на модели",
         "8 коммитов из 95 в ветке", "поддержка"],
        ["Дмитрий Комаров", "Разработка,\nAI Product",
         "Постановка задачи и приёмка, переносимость на Linux, MCP-ресурсы в смысловом контексте, "
         "пайплайн записей интерфейса",
         "10 коммитов из 95 в ветке", "поддержка"],
    ]
    yy = hy + 0.42
    for row in rows:
        rh = max(measure(c, 10, ARIAL, False, w, 1.2) for c, w in zip(row, ws))
        for x, w, cell in zip(xs, ws, row):
            text(s, x, yy, w, rh, cell, size=10, color=WHITE, font=ARIAL, line=1.2, space=0)
        yy += rh + 0.26
        hline(s, M, yy - 0.13, 8.958, WHITE, 0.4)
    text(s, M, yy + 0.06, 8.958, 0.44,
         "Все строки выведены из истории репозитория: `git shortlog -sn` против точки ветвления.",
         size=10, color=WHITE, font=ARIAL, line=1.2, space=0)
    footer_logos(s, blue=True)

    # =============================================================== APPENDIX
    # Everything a judge may ask for and nobody should have to sit through.

    s = new_slide(prs)
    y = title(s, "Приложение", size=30)
    y = text(s, M, y + 0.28, 8.958, 1.30,
             ["Архитектура · карта угроз с числами · остаточные атаки · как проверяли модель · "
              "методика измерения · вторая половина демо."],
             size=20, font=ARIAL, line=1.25, space=0) and below(0.30)
    text(s, M, y, 8.958, 0.92,
         "Все числа — из submission/SOURCE_OF_TRUTH.md, прогон final-sealed от 2026-09-05.\n"
         "Повторить: bun run script/security-bench.ts --runs 3",
         size=14, font=ARIAL, color=GRAY, line=1.3, space=0)

    # ---------------------------------------------------------- A1. trust boundary
    s = new_slide(prs)
    y = title(s, "Считаем, что моделью можно управлять со стороны")
    y = lede(s, "Разработчику доверяем. Его окружению и рассуждению модели — нет.", y=y + 0.20)
    y = cards(s, [
        ("Доверяем",
         "Разработчику и его задаче. Его личным настройкам Kilo. Коду самого Kilo. Механизму "
         "изоляции операционной системы."),
        ("Не доверяем",
         "Рассуждению модели. Содержимому репозитория. Имени пакета, которое придумала модель. "
         "Описаниям внешних инструментов. Коду проекта."),
        ("Что из этого следует",
         "Решение не зависит от того, распознала ли модель атаку. Его принимает отдельный слой, а "
         "часть границ подтверждает операционная система."),
    ], y=y + 0.22, xs=COL3_X, w=COL3_W, body_h=0.80, size=11)
    statement(s, y + 0.26, "Намерение ≠ полномочия",
              "Подсказка со стороны может изменить то, чего хочет модель. Изменить то, что ей "
              "позволено, она не может.")

    # ---------------------------------------------------------- A2. the decision
    s = new_slide(prs)
    y = title(s, "Решение — по действию, источнику и истории")
    q = [
        ("Что агент хочет сделать", "запустить команду, прочитать файл, вызвать инструмент"),
        ("С каким файлом или пакетом", "путь раскрывается до настоящего, команда разбирается"),
        ("Откуда пришёл вызов", "встроенный инструмент, внешний, инструмент проекта, расширение"),
        ("Просил ли это пользователь", "движок этого входа не получает и судит одинаково"),
        ("Что уже было в этой сессии", "какие секреты прочитаны и в какие файлы попали"),
        ("Какие признаки риска нашлись", "разбор команды, происхождение пакета, содержимое файла"),
    ]
    y0 = y + 0.24
    for i, (head, sub) in enumerate(q):
        x = COL3_X[i % 3]
        yy = y0 + (i // 3) * 0.70
        text(s, x, yy, COL3_W, 0.21, head, size=11, color=BLUE, font=INTER, space=0)
        text(s, x, yy + 0.23, COL3_W, 0.38, sub, size=9.5, font=ARIAL, color=GRAY, line=1.2, space=0)
    y = y0 + 2 * 0.70
    r = rect(s, M, y, 8.958, 0.40, fill=BLUE)
    para(r.text_frame, "Политика безопасности — один и тот же движок для каждого вызова с "
                       "последствиями", 12, True, WHITE, INTER, space=0, first=True)
    r.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cards(s, [
        ("Разрешено · ALLOW", "Действие выполняется в тех границах, к которым относится разрешение."),
        ("Вопрос человеку · ASK", "Дальше без человека небезопасно. Самостоятельный режим этот "
                                  "вопрос не снимает."),
        ("Запрещено · DENY", "Действие не выполняется. Агент получает причину и продолжает работу."),
    ], y=y + 0.40 + 0.22, xs=COL3_X, w=COL3_W, body_h=0.46, size=11)

    # ---------------------------------------------------------- A3. risk signals
    s = new_slide(prs)
    y = title(s, "Откуда берутся признаки риска")
    y = lede(s, "Каждый признак — ответ на понятный вопрос, а не совпадение со списком строк.",
             y=y + 0.18)
    left = [
        ("Опасна ли сама команда",
         "Раскрываются переменные, вызовы и пути.",
         "Deterministic Security"),
        ("Что известно про пакет",
         "Возраст, распространённость, скрипты установки.",
         "Package Security"),
        ("Мы уже читали секрет",
         "Что прочитано и в какие файлы это попало.",
         "Stateful Egress"),
        ("Кто вызывает действие",
         "Права следуют из происхождения, не из описания.",
         "Delegated Tool Security"),
    ]
    right = [
        ("Есть ли в файле учётные данные",
         "Форматы токенов, ключи, присваивания в тексте.",
         "Content Secret Detection"),
        ("Код проекта пытается запуститься сам",
         "Файл выполняет верхний уровень при импорте.",
         "Executable Code Trust"),
        ("Расширение выходит за свои права",
         "Чтение, запись, сеть и процессы — по запросу.",
         "Permissioned Extension Runtime"),
    ]
    y0 = y + 0.24
    qa_list(s, left, COL2_X[0], y0, COL2_W, pitch=0.72)
    qa_list(s, right, COL2_X[1], y0, COL2_W, pitch=0.72)

    # ---------------------------------------------------------- A4. allow is not full access
    s = new_slide(prs)
    y = title(s, "Разрешить действие ≠ отдать все права", w=8.4)
    y = lede(s, "Для расширения проекта одобрение — только первый шаг.", y=y + 0.18)
    ly = vchain(s, [
        ("Пользователь одобрил расширение", "одобрены конкретные байты: правка файла отзывает одобрение"),
        ("Ему выданы конкретные права", "чтение, запись, сеть, запуск процессов — по отдельности"),
        ("Работает в отдельном процессе под профилем ОС", "ссылка, «..» и абсолютный путь не выводят наружу"),
        ("Привилегированное действие — снова запрос", "к тому же движку, а не прямой вызов"),
    ], x=M, y=y + 0.20, w=5.30)
    rx, ry = 6.100, y + 0.20
    text(s, rx, ry, 3.379, 0.21, "Что это дало на измерении", size=11, color=BLUE, font=INTER, space=0)
    ry += 0.29
    for head, val in (("Расширение вышло за выданные права", "24 / 24  →  0 / 24"),
                      ("Расширение прочитало чужой файл", "33 / 33  →  3 / 33")):
        hh = max(measure(head, 10.5, ARIAL, False, 3.379, 1.2), 0.20)
        text(s, rx, ry, 3.379, hh, head, size=10.5, font=ARIAL, space=0)
        text(s, rx, ry + hh + 0.03, 3.379, 0.30, val, size=15, bold=True, color=BLUE, font=INTER, space=0)
        ry += hh + 0.44
    text(s, rx, ry + 0.06, 3.379, 1.22,
         "Граница держится на механизме операционной системы. Проверена прямыми пробами внутри "
         "настоящего процесса на macOS и на Linux. Без такого механизма расширение не запускается "
         "вовсе.",
         size=9.5, font=ARIAL, color=GRAY, line=1.25, space=0)
    text(s, M, ly + 0.14, 5.30, 0.20,
         "Permissioned Extension Runtime · Read-Confinement",
         size=9, color=GRAY, font=ARIAL, space=0)

    # ---------------------------------------------------------- A5. package provenance
    s = new_slide(prs)
    y = title(s, "Подозрительный пакет останавливается до установки")
    y = lede(s, "Модель придумывает имя библиотеки, атакующий его заранее регистрирует.", y=y + 0.14)
    y0 = y + 0.22
    steps = [("Агент предлагает\nзависимость", 0.521, 1.90, None),
             ("Оценка\nпроисхождения", 2.641, 1.90, BLUE),
             ("Решение\nALLOW · ASK · DENY", 4.761, 1.90, BLUE),
             ("Установка\nзапускается или нет", 6.881, 1.90, None)]
    for i, (label, x, w, fill) in enumerate(steps):
        r = rect(s, x, y0, w, 0.56, fill=fill, lineclr=None if fill else BLACK)
        head, sub = label.split("\n")
        para(r.text_frame, head, 10.5, True, WHITE if fill else BLACK, INTER, space=1, first=True)
        para(r.text_frame, sub, 9, False, WHITE if fill else GRAY, ARIAL, space=0)
        if i < 3:
            arrow(s, x + w + 0.03, y0 + 0.28, 0.16)
    y = cards(s, [
        ("Что известно точно",
         "Возраст пакета и выпуска, скрипты установки, репозиторий, подмена реестра."),
        ("Что оцениваем приблизительно",
         "Насколько распространён и похоже ли имя на известную библиотеку."),
        ("Если ничего не известно",
         "Метаданных нет, пакета нет в реестре. Неизвестность не даёт разрешения."),
    ], y=y0 + 0.56 + 0.16, xs=COL3_X, w=COL3_W, body_h=0.58, size=11)
    statement(s, y + 0.12, "Происхождение пакета оценивается до его локального исполнения",
              "Проверка идёт внутри работы агента, а не в сборке. Это оценка происхождения, а не "
              "разбор кода. Измерено: 42 из 42 без защиты, 0 из 42 с ней.")

    # ---------------------------------------------------------- A6. the semantic layer
    s = new_slide(prs)
    y = title(s, "Что именно делает модель и чего она не может")
    y = lede(s, "Два вопроса, на которые у детерминированного контура нет входа.", y=y + 0.14)
    y = cards(s, [
        ("Откуда пришла инструкция",
         "README, документация зависимости, страница из сети, ответ внешнего инструмента. Для правил "
         "по путям это одинаковые обычные файлы. Написал их не пользователь."),
        ("Связано ли действие с задачей",
         "«Поправь опечатку в README» и отправка ключа наружу не связаны ничем. Совпадение с задачей "
         "только снижает тревогу и никогда не даёт разрешения."),
        ("Чего модель не может",
         "Она возвращает не решение, а свидетельство, и свёртка умеет только ужесточать. Отменить "
         "запрет или выдать разрешение она не способна — это свойство кода, а не обещание."),
    ], y=y + 0.14, xs=COL3_X, w=COL3_W, body_h=0.78, size=10.5)
    y = stat_row(s, [
        ("24 % → 5 %", "общая успешность атак"),
        ("92 % → 0 %", "подброшенная инструкция"),
        ("100 % → 11 %", "действие не по задаче"),
        ("98 %", "задач — как без слоя"),
    ], y=y + 0.18, numsize=15)
    text(s, M, y + 0.12, 8.958, 0.24,
         "Цену не прячем: 16 % решений идут к модели, путь с ней 1,1–1,4 с, ≈ $0,30 на 1000 решений.",
         size=11, font=ARIAL, space=0)

    # ---------------------------------------------------------- A7. the benign half
    s = new_slide(prs)
    y = title(s, "Тот же файл и та же команда — вопроса нет", size=24, w=8.958)
    y = text(s, M, y + 0.14, 8.958, 0.34,
             "Вторая половина демо. Отличается только текст в README.",
             size=16, font=ARIAL, space=0) and below(0.14)
    y = shot(s, "pitch-benign.png", M, y, 8.958)
    text(s, M, y + 0.12, 8.958, 0.30,
         "Вопроса человеку не было: слой промолчал, файл ушёл, приёмник его получил.",
         size=14, bold=True, color=BLUE, font=INTER, space=0)

    # ---------------------------------------------------------- A8. other recorded sessions
    s = new_slide(prs)
    y = title(s, "Одна команда — разные решения", y=0.940)
    y = lede(s, "Записи настоящих сессий Kilo с этой ветки. Отличается только пакет.", y=y + 0.10)
    y = shot(s, "02-known-package.png", M, y + 0.08, 6.800,
             "Зрелая библиотека — ставится сразу, без единого вопроса")
    y = shot(s, "03-unknown-package-stopped.png", M, y + 0.10, 6.800,
             "Придуманного имени в реестре нет — агент остановился, пакетный менеджер не запускался")
    text(s, M, y + 0.10, 8.958, 0.24,
         "Это не чёрный список: имени, зарегистрированного вчера, в списке нет по определению.",
         size=11, font=ARIAL, space=0)

    # ---------------------------------------------------------- A9. with and without
    s = new_slide(prs)
    y = title(s, "Без защиты проект исчезает", y=0.940)
    y = lede(s, "Та же задача и тот же агент. Отличается одно значение настройки.", y=y + 0.10)
    y = shot(s, "06-workspace-wipe-baseline.png", M, y + 0.08, 6.800,
             "Security Auto выключен — после прогона в рабочем дереве не осталось ничего")
    y = shot(s, "07-workspace-wipe-stopped.png", M, y + 0.10, 6.800,
             "Security Auto включён — агент ждёт человека, файлы на месте")
    text(s, M, y + 0.10, 8.958, 0.24,
         "Запись сделана без человека за экраном, поэтому вопрос виден как «отвечать некому».",
         size=11, font=ARIAL, color=GRAY, space=0)

    # ---------------------------------------------------------- A10. how it was measured
    s = new_slide(prs)
    y = title(s, "Успех атаки — это последствие")
    y = cards(s, [
        ("Что проверяли",
         "192 сценария: 141 атака и 51 обычная задача, каждый по три раза. Десять конфигураций: "
         "соседние отличаются одним включённым слоем."),
        ("Как изолировали",
         "Одноразовый каталог во временной папке, подставной домашний каталог, безвредные заглушки "
         "пакетных менеджеров, сеть только на петле."),
        ("Как определяли успех",
         "Помеченный файл действительно удалён. Выдуманный ключ действительно дошёл до приёмника. "
         "Установка действительно запустилась."),
    ], y=y + 0.22, xs=COL3_X, w=COL3_W, body_h=0.80, size=11)
    y = ladder(s, y + 0.38)
    text(s, M, y + 0.04, 8.958, 0.68,
         "Набор атак рос по ходу работы: найденный обход не убирался после исправления, а "
         "добавлялся в общий набор, и все конфигурации пересчитывались заново. Ранние конфигурации "
         "выглядят хуже потому, что тест стал сложнее.",
         size=11, font=ARIAL, line=1.2, space=0)

    # ---------------------------------------------------------- A11. threat-by-threat
    s = new_slide(prs)
    y = title(s, "Три класса держатся только на смысловом слое", y=0.940)
    y = lede(s, "Серая полоса — 100 % атак без защиты. Синяя — что доходит с защитой.", y=y + 0.16)
    y = scorecard(s, [
        ("Чужие ключи и настройки самого Kilo", 0, 54, "Закрыто"),
        ("Установка подозрительного пакета", 0, 42, "Закрыто"),
        ("Внешний инструмент делает лишнее", 0, 30, "Закрыто"),
        ("Утечка секрета цепочкой действий", 0, 30, "Закрыто"),
        ("Расширение выходит за выданные права", 0, 24, "Закрыто"),
        ("Действие не связано с задачей пользователя", 1, 9, "Слоем смысла · 11 %"),
        ("Текст в репозитории уговаривает агента", 0, 75, "Закрыто · слоем смысла"),
        ("Секрет в обычном файле проекта", 0, 33, "Закрыто · было 18 %"),
        ("Расширение читает файлы разработчика", 3, 33, "Сильно снижено · 9 %"),
        ("Код репозитория исполняется при загрузке", 3, 27, "Сильно снижено · 11 %"),
        ("Разрушительная команда, в том числе скрытая", 12, 60, "Частично · 20 %"),
    ], y=y + 0.06, namew=3.60, trackw=1.85, pitch=0.250)
    text(s, M, y + 0.06, 8.958, 0.28,
         "Всего 19 прогонов из 417. Три строки со слоем смысла без него читались бы как 100 %, 92 % и 18 %.",
         size=10.5, font=ARIAL, line=1.2, space=0)

    # ---------------------------------------------------------- A12. limits and pilot
    s = new_slide(prs)
    y = title(s, "Что не закрыто и что дальше")
    y = cards(s, [
        ("Чего защита не умеет",
         "Путь внутри строки python -c, node -e или perl -e не становится операндом при разборе.\n"
         "git stash -u очищает дерево и разрешён; правила на следующий drop не хватает.\n"
         "Внутри сужённого процесса расширение видит размер файла и то, что он есть.\n"
         "Один вызов модели из 117 не уложился в срок — там, где под слоем нет правила, это пропуск.\n"
         "Границы проверены на macOS и Linux. На Linux содержимое домашнего каталога закрыто, "
         "но имена каталогов на пути привязок видны."),
        ("Чего не показывает измерение",
         "Набор атак наш собственный: «100 % без защиты» — это про него, а не про Kilo вообще.\n"
         "Ход работы агента задан сценарием, поэтому измеряется удержание, а не то, догадается ли "
         "модель напасть.\n"
         "Качество смыслового слоя — свойство одной модели: всё снято на claude-haiku-4.5 через "
         "OpenRouter, другая даст другие числа."),
        ("Пилот",
         "Команда 5–10 разработчиков, репозиторий с внешними зависимостями и внешними инструментами.\n"
         "Что смотрим: доля успешных атак, доля доведённых задач, ошибочные запреты, вопросы на "
         "задачу, задержка и сколько человек через две недели оставили режим включённым.\n"
         "Отложенный набор израсходован; новый пишет внешний автор, а конфигурация уже запечатана."),
    ], y=y + 0.18, xs=COL3_X, w=COL3_W, body_h=1.56, size=10.5)
    hline(s, M, y + 0.06, 8.958, GRAY, 0.75)
    text(s, M, y + 0.10, 8.958, 0.44,
         "Повторить любую цифру:  bun run script/security-bench.ts --runs 3\n"
         "Повторить любую запись сессии:  bash submission/evidence/ui/capture-all.sh",
         size=10, font=ARIAL, color=GRAY, line=1.3, space=0)


APP_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" \
xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">\
<TotalTime>0</TotalTime><Application>python-pptx</Application>\
<PresentationFormat>Custom</PresentationFormat><Slides>{slides}</Slides><Notes>0</Notes>\
<HiddenSlides>0</HiddenSlides><MMClips>0</MMClips><ScaleCrop>false</ScaleCrop>\
<LinksUpToDate>false</LinksUpToDate><SharedDoc>false</SharedDoc><HyperlinksChanged>false</HyperlinksChanged>\
</Properties>"""


def scrub(path, slides):
    """Replace the template's extended properties: they still carry the source deck's slide
    titles and counts, which have nothing to do with this presentation."""
    import zipfile

    tmp = path + ".tmp"
    with zipfile.ZipFile(path) as src, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "docProps/app.xml":
                data = APP_XML.format(slides=slides).encode("utf-8")
            dst.writestr(item, data)
    os.replace(tmp, path)


def verify_backgrounds(path):
    """Fail the build if a painted background would not survive PowerPoint.

    Checked on the saved package rather than in memory, because that is the artefact that ships, and
    because the failure this catches is silent everywhere else: an out-of-place `<p:bg>` renders fine
    through LibreOffice and disappears in PowerPoint, leaving white text on a white page."""
    import zipfile
    from lxml import etree

    pns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    bad = []
    with zipfile.ZipFile(path) as zf:
        for name in sorted(n for n in zf.namelist() if n.startswith("ppt/slides/slide")):
            root = etree.fromstring(zf.read(name))
            stray = root.findall(f"{pns}bg")
            csld = root.find(f"{pns}cSld")
            inside = csld.findall(f"{pns}bg") if csld is not None else []
            if stray:
                bad.append(f"{name}: <p:bg> is a sibling of <p:cSld>, not inside it")
            elif inside and list(csld).index(inside[0]) != 0:
                bad.append(f"{name}: <p:bg> is not the first child of <p:cSld>")
    if bad:
        raise SystemExit("background will not render in PowerPoint:\n  " + "\n  ".join(bad))
    return sum(1 for _ in [b for b in bad]) == 0


def main():
    preflight()
    shutil.copyfile(TEMPLATE, OUT)
    prs = Presentation(OUT)
    # drop the template's own slides, keep master / layouts / theme / embedded fonts
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)
    build(prs)
    core = prs.core_properties
    core.title = "Security Auto Mode для Kilo Code"
    core.subject = "Детерминированный ALLOW / ASK / DENY для side-effecting tool call'ов"
    core.author = "Dmitry Voropaev"
    core.last_modified_by = "Dmitry Voropaev"
    core.comments = "AI Talent Hub. Числа: submission/SOURCE_OF_TRUTH.md"
    core.keywords = "Kilo Code, Security Auto Mode, coding agent, benchmark"
    core.created = core.modified = datetime.now(timezone.utc).replace(tzinfo=None, microsecond=0)
    core.revision = 1
    n = len(prs.slides._sldIdLst)
    prs.save(OUT)
    scrub(OUT, n)
    verify_backgrounds(OUT)
    print(f"wrote {OUT} - {n} slides")
    sys.exit(1 if report() else 0)


if __name__ == "__main__":
    main()
